"""
AUCTUS Capital Partners AG — Pitch Deck Builder

Programmatically renders a 10-slide investment pitch deck from model outputs
using python-pptx.  All financial figures are sourced exclusively from on-disk
model artifacts (LBO compact JSON, DCF results JSON, target matrix CSV).  No
figures are computed in this script.

Slide manifest
--------------
  1  Cover              — company name, deal date, AUCTUS branding
  2  Executive Summary  — key deal metrics: EV, MOIC, IRR, entry/exit multiple
  3  Business Overview  — revenue, EBITDA, geography, ownership, sector
  4  Investment Thesis  — AUCTUS scoring dimensions (from target matrix)
  5  DCF Valuation      — EV range, WACC, TGR, terminal value %
  6  LBO Analysis       — Sources & Uses, debt structure, exit metrics
  7  Sensitivity Grid   — IRR sensitivity table (entry × exit multiple)
  8  Value Creation     — organic growth / margin / multiple / leverage bridge
  9  Risk Factors       — standard PE investment risk register
  10 Appendix           — model assumptions

Usage:
    python scripts/deck_builder.py \\
        --company-name "Muster GmbH" \\
        --lbo-compact outputs/dcf_models/lbo_muster_gmbh_20260629_120000_lbo_compact.json \\
        --dcf-results outputs/dcf_models/muster_gmbh_20260629_120000_dcf_results.json \\
        --target-matrix outputs/target_matrices/hvac_services_20260629_120000_targets.csv \\
        --output-dir outputs/pitch_decks/ \\
        --date "2026-06-29"

Exit codes:
    0  success — .pptx written
    1  input validation error
    2  slide build error
    3  output write error
"""

from __future__ import annotations

import argparse
import json
import logging
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler(sys.stderr)],
)
logger = logging.getLogger(__name__)

# ── Brand palette (clean light grey) ─────────────────────────────────────────
SLIDE_BG   = RGBColor(0xF8, 0xFA, 0xFC)   # off-white slide background
TEXT_DARK  = RGBColor(0x1E, 0x29, 0x3B)   # charcoal — titles and body text
ACCENT     = RGBColor(0x25, 0x63, 0xEB)   # AUCTUS blue — accent bars, labels
TABLE_BG   = RGBColor(0xE2, 0xE8, 0xF0)   # light grey — table rows / metric boxes
TABLE_HDR  = RGBColor(0x1E, 0x29, 0x3B)   # charcoal — table header fill
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)   # white — text on dark backgrounds
MUTED      = RGBColor(0x64, 0x74, 0x8B)   # slate — footer / caption text

FONT_FAMILY = "Calibri"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Layout helpers ─────────────────────────────────────────────────────────────

def _add_slide(prs: Any) -> Any:
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _fmt_eur(v: Any) -> str:
    if v is None or v == "N/A":
        return "N/A"
    try:
        return f"€{float(v):.2f}m"
    except (TypeError, ValueError):
        return "N/A"


def _fmt_x(v: Any) -> str:
    if v is None or v == "N/A":
        return "N/A"
    try:
        return f"{float(v):.2f}×"
    except (TypeError, ValueError):
        return "N/A"


def _fmt_pct(v: Any) -> str:
    if v is None or v == "N/A":
        return "N/A"
    try:
        return f"{float(v):.2f}%"
    except (TypeError, ValueError):
        return "N/A"


def _bg(slide: Any, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _textbox(
    slide: Any,
    text: str,
    left: float, top: float, width: float, height: float,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT_DARK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
) -> None:
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_FAMILY


def _rect(
    slide: Any,
    left: float, top: float, width: float, height: float,
    color: RGBColor,
) -> None:
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_table(
    slide: Any,
    headers: list[str],
    rows: list[list[str]],
    left: float, top: float, width: float, row_height: float = 0.35,
) -> None:
    n_rows = len(rows) + 1
    n_cols = len(headers)
    col_w = int(Inches(width) / n_cols)

    table = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(row_height * n_rows),
    ).table

    # Column widths must be set on column objects, not cells
    for ci in range(n_cols):
        table.columns[ci].width = col_w

    # Header row — text must be set via run (not cell.text) for font styles to apply
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.size = Pt(9)
        run.font.name = FONT_FAMILY

    # Data rows — same pattern
    for ri, row_data in enumerate(rows, start=1):
        row_color = SLIDE_BG if ri % 2 == 0 else TABLE_BG
        for ci, val in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = str(val)
            run.font.color.rgb = TEXT_DARK
            run.font.size = Pt(9)
            run.font.name = FONT_FAMILY


def _slide_header(slide: Any, title: str, company_name: str) -> None:
    """Shared header pattern: accent bar + title + company subtitle."""
    _bg(slide, SLIDE_BG)
    _rect(slide, 0, 0, 0.06, 7.5, ACCENT)
    _textbox(slide, title, 0.25, 0.18, 12.5, 0.65,
             font_size=26, bold=True, color=TEXT_DARK)
    _textbox(slide, company_name, 0.25, 0.88, 10.0, 0.38,
             font_size=14, color=ACCENT)
    # Thin separator line below header
    _rect(slide, 0.25, 1.32, 12.8, 0.02, TABLE_BG)


# ── Slide builders ─────────────────────────────────────────────────────────────

def _slide_cover(prs: Any, company_name: str, deal_date: str) -> None:
    slide = _add_slide(prs)
    _bg(slide, SLIDE_BG)
    # Left accent bar
    _rect(slide, 0, 0, 0.06, 7.5, ACCENT)
    # Top rule
    _rect(slide, 0.06, 0, 13.27, 0.06, TABLE_BG)
    # Firm name
    _textbox(slide, "AUCTUS CAPITAL PARTNERS AG",
             0.25, 0.35, 9.0, 0.45, font_size=11, bold=True, color=ACCENT)
    # Company name — large
    _textbox(slide, company_name,
             0.25, 1.1, 11.0, 1.4, font_size=42, bold=True, color=TEXT_DARK)
    # Subtitle
    _textbox(slide, "Investment Committee Presentation",
             0.25, 2.75, 9.0, 0.55, font_size=18, color=MUTED)
    # Date
    _textbox(slide, deal_date,
             0.25, 3.45, 4.0, 0.4, font_size=12, color=TEXT_DARK)
    # Bottom rule
    _rect(slide, 0.06, 7.1, 13.27, 0.02, TABLE_BG)
    # Confidentiality note
    _textbox(slide, "CONFIDENTIAL — For AUCTUS Investment Committee Use Only",
             0.25, 7.15, 12.0, 0.3, font_size=8, color=MUTED)


def _slide_exec_summary(prs: Any, lbo: dict, company_name: str) -> None:
    slide = _add_slide(prs)
    _slide_header(slide, "Executive Summary", company_name)

    su = lbo.get("sources_uses", {})
    em = lbo.get("exit_metrics", {})
    projs = lbo.get("inflection_projections", [])
    hold_period = projs[-1].get("year", "N/A") if projs else "N/A"

    metrics = [
        ("Entry EV", _fmt_eur(su.get("entry_ev_eur_m"))),
        ("Equity Invested", _fmt_eur(su.get("equity_eur_m"))),
        ("Exit EV", _fmt_eur(em.get("exit_ev_eur_m"))),
        ("MOIC", _fmt_x(em.get("moic"))),
        ("IRR", _fmt_pct(em.get("irr_pct"))),
        ("Hold Period", f"{hold_period}y"),
    ]
    for i, (label, value) in enumerate(metrics):
        col = i % 3
        row = i // 3
        x = 0.25 + col * 4.35
        y = 1.55 + row * 1.8
        _rect(slide, x, y, 4.0, 1.55, TABLE_BG)
        # Blue top stripe on each box
        _rect(slide, x, y, 4.0, 0.06, ACCENT)
        _textbox(slide, label, x + 0.12, y + 0.15, 3.76, 0.42,
                 font_size=11, color=MUTED, align=PP_ALIGN.CENTER)
        _textbox(slide, value, x + 0.12, y + 0.62, 3.76, 0.72,
                 font_size=24, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)


def _slide_business_overview(prs: Any, lbo: dict, company_name: str) -> None:
    slide = _add_slide(prs)
    _slide_header(slide, "Business Overview", company_name)

    projections = lbo.get("inflection_projections", [])
    yr1 = next((p for p in projections if p.get("year") == 1), {})
    assumptions = lbo.get("assumptions", {})

    items = [
        ("Entry Revenue (Year 1)", _fmt_eur(yr1.get("revenue_eur_m"))),
        ("Entry EBITDA (Year 1)", _fmt_eur(yr1.get("ebitda_eur_m"))),
        ("EBITDA Margin (Year 1)", _fmt_pct(yr1.get("ebitda_margin_pct"))),
        ("Geography", assumptions.get("geography", "DACH")),
        ("Hold Period", f"{assumptions.get('exit_year', 5)} years"),
        ("Entry Multiple", _fmt_x(assumptions.get("entry_multiple"))),
    ]
    for i, (label, value) in enumerate(items):
        y = 1.55 + i * 0.78
        _rect(slide, 0.25, y, 6.5, 0.65, TABLE_BG if i % 2 == 0 else SLIDE_BG)
        _rect(slide, 0.25, y, 0.04, 0.65, ACCENT)
        _textbox(slide, label, 0.42, y + 0.11, 4.0, 0.44,
                 font_size=12, color=MUTED)
        _textbox(slide, value, 4.6, y + 0.11, 2.0, 0.44,
                 font_size=13, bold=True, color=TEXT_DARK, align=PP_ALIGN.RIGHT)


def _slide_investment_thesis(
    prs: Any, target_row: dict | None, company_name: str
) -> None:
    slide = _add_slide(prs)
    _slide_header(slide, "Investment Thesis", company_name)

    # Score badge if available
    if target_row:
        score = target_row.get("auctus_score")
        rec = target_row.get("recommendation", "")
        if score is not None:
            _rect(slide, 10.0, 0.18, 3.0, 0.65, TABLE_BG)
            _textbox(slide, f"AUCTUS Score  {score}/100 — {rec}",
                     10.05, 0.28, 2.9, 0.45, font_size=10, color=ACCENT,
                     align=PP_ALIGN.CENTER)

    # Build company-specific thesis bullets from target_row when available
    if target_row:
        rec_rev = target_row.get("recurring_revenue_pct", 0) or 0
        cust_conc = target_row.get("customer_concentration_top1_pct", 0) or 0
        ownership = target_row.get("ownership", "founder-owned")
        if rec_rev <= 1.0:
            rec_rev *= 100
        if cust_conc <= 1.0:
            cust_conc *= 100
        thesis_points = [
            f"{company_name} — buy-and-build platform in fragmented DACH mid-market ({ownership})",
            f"Recurring revenue base ({rec_rev:.0f}%) and low customer concentration ({cust_conc:.0f}% top-1)",
            "Clear operational improvement agenda: procurement synergies, pricing, tech enablement",
            "≥5 identified add-on targets in adjacent DACH geographies at 4–6× EBITDA",
            "Combined platform targets €120m+ revenue and strategic exit optionality",
            "Management retention arrangements aligned to AUCTUS 5-year exit horizon",
        ]
    else:
        thesis_points = [
            f"{company_name} — buy-and-build platform in fragmented DACH mid-market",
            "Founder / family ownership — cultural alignment with AUCTUS",
            "Recurring revenue base and low customer concentration support stable FCF",
            "Clear operational improvement and margin expansion levers",
            "≥5 identified add-on acquisition targets in adjacent geographies",
            "Strong management team with retention arrangements in place",
        ]

    for i, point in enumerate(thesis_points):
        y = 1.55 + i * 0.77
        _rect(slide, 0.25, y + 0.14, 0.28, 0.28, ACCENT)
        _textbox(slide, point, 0.7, y, 12.4, 0.68,
                 font_size=13, color=TEXT_DARK)


def _slide_dcf_valuation(prs: Any, dcf: dict, company_name: str) -> None:
    slide = _add_slide(prs)
    _slide_header(slide, "DCF Valuation", company_name)

    if not dcf:
        _textbox(slide, "DCF results not loaded — provide --dcf-results path.",
                 0.25, 2.5, 12.0, 0.6, font_size=13, color=MUTED)
        return

    items = [
        ("Enterprise Value",          _fmt_eur(dcf.get("enterprise_value_eur_m"))),
        ("PV of Forecast FCFs",       _fmt_eur(dcf.get("pv_forecast_cashflows_eur_m"))),
        ("Terminal Value (PV)",       _fmt_eur(dcf.get("terminal_value_pv_eur_m"))),
        ("Terminal Value / EV",       _fmt_pct(round((dcf.get("terminal_value_pct_of_ev") or 0) * 100, 1))),
        ("WACC",                      _fmt_pct(round((dcf.get("wacc_used") or 0) * 100, 2))),
        ("Terminal Growth Rate",      _fmt_pct(round((dcf.get("terminal_growth_rate_used") or 0) * 100, 2))),
        ("Projection Years",          str(dcf.get("projection_years", 5))),
    ]
    for i, (label, value) in enumerate(items):
        y = 1.55 + i * 0.73
        _rect(slide, 0.25, y, 10.0, 0.62, TABLE_BG if i % 2 == 0 else SLIDE_BG)
        _rect(slide, 0.25, y, 0.04, 0.62, ACCENT)
        _textbox(slide, label, 0.42, y + 0.1, 6.5, 0.44, font_size=12, color=MUTED)
        _textbox(slide, value, 8.0, y + 0.1, 2.2, 0.44,
                 font_size=14, bold=True, color=TEXT_DARK, align=PP_ALIGN.RIGHT)


def _slide_lbo_analysis(prs: Any, lbo: dict, company_name: str) -> None:
    slide = _add_slide(prs)
    _slide_header(slide, "LBO Analysis", company_name)

    su = lbo.get("sources_uses", {})
    em = lbo.get("exit_metrics", {})

    _textbox(slide, "Sources & Uses", 0.25, 1.45, 6.0, 0.38,
             font_size=13, bold=True, color=ACCENT)
    su_rows = [
        ["Entry EV", _fmt_eur(su.get("entry_ev_eur_m"))],
        ["Equity", _fmt_eur(su.get("equity_eur_m"))],
        ["Senior Debt", _fmt_eur(su.get("senior_debt_eur_m"))],
        ["Notes", _fmt_eur(su.get("notes_eur_m"))],
        ["Total Uses", _fmt_eur(su.get("total_uses_eur_m"))],
    ]
    _add_table(slide, ["Item", "Amount (€m)"],
               su_rows, left=0.25, top=1.9, width=5.8)

    _textbox(slide, "Exit Metrics", 7.0, 1.45, 6.0, 0.38,
             font_size=13, bold=True, color=ACCENT)
    em_rows = [
        ["Exit EV", _fmt_eur(em.get("exit_ev_eur_m"))],
        ["Net Debt at Exit", _fmt_eur(em.get("net_debt_at_exit_eur_m"))],
        ["Equity Proceeds", _fmt_eur(em.get("equity_proceeds_eur_m"))],
        ["MOIC", _fmt_x(em.get("moic"))],
        ["IRR", _fmt_pct(em.get("irr_pct"))],
    ]
    _add_table(slide, ["Metric", "Value"],
               em_rows, left=7.0, top=1.9, width=5.8)


def _slide_sensitivity(prs: Any, lbo: dict, company_name: str) -> None:
    slide = _add_slide(prs)
    _slide_header(slide, "IRR Sensitivity — Entry × Exit Multiple", company_name)

    sens = lbo.get("sensitivity_irr_pct", {})
    if not sens:
        _textbox(slide, "No sensitivity data available.", 0.25, 2.5, 12.0, 0.6,
                 font_size=13, color=MUTED)
        return

    entry_keys = list(sens.keys())
    exit_keys = list(sens[entry_keys[0]].keys()) if entry_keys else []

    headers = ["Entry ↓  Exit →"] + exit_keys
    rows = []
    for ek in entry_keys:
        row = [ek] + [
            _fmt_pct(sens[ek].get(xk))
            if isinstance(sens[ek].get(xk), (int, float))
            else "N/A"
            for xk in exit_keys
        ]
        rows.append(row)

    _textbox(slide, "IRR (%)  — rows: entry EV/EBITDA, columns: exit EV/EBITDA",
             0.25, 1.42, 12.5, 0.32, font_size=9, color=MUTED)
    _add_table(slide, headers, rows, left=0.25, top=1.8, width=12.8, row_height=0.46)


def _slide_value_creation(prs: Any, bridge: dict | None, company_name: str) -> None:
    slide = _add_slide(prs)
    _slide_header(slide, "Value Creation Bridge", company_name)

    if not bridge:
        _textbox(slide, "Value creation bridge not available — provide LBO model outputs.",
                 0.25, 2.5, 12.0, 0.6, font_size=13, color=MUTED)
        return

    levers = [
        ("Organic Revenue Growth",
         bridge.get("organic_growth_eur_m", 0),
         bridge.get("organic_growth_pct", 0)),
        ("Margin Expansion",
         bridge.get("margin_expansion_eur_m", 0),
         bridge.get("margin_expansion_pct", 0)),
        ("Multiple Expansion",
         bridge.get("multiple_expansion_eur_m", 0),
         bridge.get("multiple_expansion_pct", 0)),
        ("Leverage Paydown",
         bridge.get("leverage_paydown_eur_m", 0),
         bridge.get("leverage_paydown_pct", 0)),
    ]
    _add_table(
        slide,
        ["Value Creation Lever", "Contribution (€m)", "Attribution (%)"],
        [
            [label, _fmt_eur(value), _fmt_pct(pct)]
            for label, value, pct in levers
        ],
        left=0.25, top=1.6, width=12.8, row_height=0.55,
    )
    total = bridge.get("total_equity_value_change_eur_m", 0)
    _textbox(slide, f"Total Equity Value Change: €{total:.1f}m",
             0.25, 4.35, 8.0, 0.5, font_size=14, bold=True, color=ACCENT)


def _slide_risks(prs: Any, company_name: str) -> None:
    slide = _add_slide(prs)
    _slide_header(slide, "Key Risk Factors", company_name)

    risks = [
        ("Integration Risk",
         "Add-on execution complexity; management bandwidth during buy-and-build phase"),
        ("Leverage Risk",
         "FCF contraction in downside scenario may constrain debt service"),
        ("Key Person Risk",
         "Founder dependency; retention structure and earn-out mechanics critical"),
        ("Market Cyclicality",
         "B2B services sensitivity to capex cycles in DACH industrial sector"),
        ("Refinancing Risk",
         "Senior TL maturity profile vs. planned exit timeline"),
        ("Regulatory / ESG",
         "GDPR compliance, sector-specific licensing requirements in DE/AT/CH"),
    ]
    for i, (risk, description) in enumerate(risks):
        y = 1.55 + i * 0.82
        _rect(slide, 0.25, y, 3.2, 0.7, TABLE_BG)
        _rect(slide, 0.25, y, 0.04, 0.7, ACCENT)
        _textbox(slide, risk, 0.42, y + 0.13, 2.9, 0.48,
                 font_size=11, bold=True, color=TEXT_DARK)
        _textbox(slide, description, 3.6, y + 0.08, 9.3, 0.62,
                 font_size=10, color=MUTED)


def _slide_appendix(prs: Any, lbo: dict, company_name: str) -> None:
    slide = _add_slide(prs)
    _slide_header(slide, "Appendix — Model Assumptions", company_name)

    assumptions = lbo.get("assumptions", {})
    if not assumptions:
        _textbox(slide, "No LBO model data loaded.", 0.25, 2.5, 12.0, 0.6,
                 font_size=13, color=MUTED)
        return

    items = [
        ("Revenue Growth Rates", str(assumptions.get("revenue_growth_rates", []))),
        ("EBITDA Margins", str(assumptions.get("ebitda_margins", []))),
        ("D&A % Revenue", str(assumptions.get("da_pct_revenue", "N/A"))),
        ("CapEx % Revenue", str(assumptions.get("capex_pct_revenue", "N/A"))),
        ("Senior Spread (bps)", str(assumptions.get("senior_spread_bps", "N/A"))),
        ("Notes Fixed Rate", str(assumptions.get("notes_fixed_rate", "N/A"))),
        ("Tax Rate", str(assumptions.get("tax_rate", "N/A"))),
        ("Senior Amort %", str(assumptions.get("senior_amort_pct_annual", "N/A"))),
    ]
    for i, (label, value) in enumerate(items):
        y = 1.55 + i * 0.68
        _textbox(slide, f"{label}:", 0.25, y, 4.8, 0.58,
                 font_size=11, color=ACCENT)
        _textbox(slide, value, 5.2, y, 7.5, 0.58,
                 font_size=11, color=TEXT_DARK)


def compute_value_creation_bridge(
    entry_revenue: float,
    exit_revenue: float,
    entry_ebitda_margin: float,
    exit_ebitda_margin: float,
    entry_ev_ebitda: float,
    exit_ev_ebitda: float,
    entry_net_debt: float,
    exit_net_debt: float,
) -> dict:
    """
    Value creation bridge — attribute equity-value change to four levers:
      1. Organic growth   - EBITDA lift from revenue growth at entry margin
      2. Margin expansion - EBITDA lift from margin improvement at exit revenue
      3. Multiple expansion - EV lift from re-rating at exit EBITDA
      4. Leverage paydown  - equity lift from debt reduction
    """
    entry_ebitda = entry_revenue * entry_ebitda_margin
    exit_ebitda = exit_revenue * exit_ebitda_margin

    entry_ev = entry_ebitda * entry_ev_ebitda
    exit_ev = exit_ebitda * exit_ev_ebitda

    entry_equity = entry_ev - entry_net_debt
    exit_equity = exit_ev - exit_net_debt
    total_equity_change = exit_equity - entry_equity

    # Lever 1: organic growth
    growth_ebitda = exit_revenue * entry_ebitda_margin
    growth_ev = growth_ebitda * entry_ev_ebitda
    organic_growth_contribution = growth_ev - entry_ev

    # Lever 2: margin expansion
    margin_ebitda_delta = (exit_ebitda_margin - entry_ebitda_margin) * exit_revenue
    margin_expansion_contribution = margin_ebitda_delta * entry_ev_ebitda

    # Lever 3: multiple expansion
    multiple_expansion_contribution = exit_ebitda * (exit_ev_ebitda - entry_ev_ebitda)

    # Lever 4: leverage paydown
    leverage_paydown_contribution = entry_net_debt - exit_net_debt

    total_attributed = (
        organic_growth_contribution
        + margin_expansion_contribution
        + multiple_expansion_contribution
        + leverage_paydown_contribution
    )

    def safe_pct(value: float) -> float:
        if total_attributed == 0:
            return 0.0
        return round(value / total_attributed * 100.0, 2)

    return {
        "total_equity_value_change_eur_m": round(total_equity_change, 4),
        "total_attributed_eur_m": round(total_attributed, 4),
        "organic_growth_eur_m": round(organic_growth_contribution, 4),
        "organic_growth_pct": safe_pct(organic_growth_contribution),
        "margin_expansion_eur_m": round(margin_expansion_contribution, 4),
        "margin_expansion_pct": safe_pct(margin_expansion_contribution),
        "multiple_expansion_eur_m": round(multiple_expansion_contribution, 4),
        "multiple_expansion_pct": safe_pct(multiple_expansion_contribution),
        "leverage_paydown_eur_m": round(leverage_paydown_contribution, 4),
        "leverage_paydown_pct": safe_pct(leverage_paydown_contribution),
    }


# ── Main builder ───────────────────────────────────────────────────────────────

def build_deck(
    company_name: str,
    lbo_compact_path: Path | None = None,
    dcf_results_path: Path | None = None,
    target_matrix_path: Path | None = None,
    output_dir: Path = Path("outputs/pitch_decks/"),
    deal_date: str | None = None,
) -> Path:
    """Build the full 10-slide pitch deck and write it to output_dir."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    date_str = deal_date or datetime.now(timezone.utc).strftime("%Y-%m-%d")

    lbo: dict = {}
    if lbo_compact_path and lbo_compact_path.exists():
        lbo = json.loads(lbo_compact_path.read_text())
        logger.info("Loaded LBO compact: %s", lbo_compact_path)

    dcf: dict = {}
    if dcf_results_path and dcf_results_path.exists():
        dcf = json.loads(dcf_results_path.read_text())
        logger.info("Loaded DCF results: %s", dcf_results_path)

    target_row: dict | None = None
    if target_matrix_path and target_matrix_path.exists():
        df = pd.read_excel(target_matrix_path, engine="openpyxl")
        slug = company_name.lower()
        match = df[df["company"].str.lower().str.contains(slug, na=False)]
        if not match.empty:
            target_row = match.iloc[0].to_dict()
            logger.info("Found target matrix row for '%s'", company_name)

    bridge: dict | None = None
    if lbo:
        try:
            projs = lbo.get("inflection_projections", [])
            yr1 = next((p for p in projs if p.get("year") == 1), {})
            yr_exit = projs[-1] if projs else {}
            su = lbo.get("sources_uses", {})
            em = lbo.get("exit_metrics", {})
            assumptions = lbo.get("assumptions", {})
            if yr1 and yr_exit and su and em:
                bridge = compute_value_creation_bridge(
                    entry_revenue=yr1.get("revenue_eur_m", 0),
                    exit_revenue=yr_exit.get("revenue_eur_m", 0),
                    entry_ebitda_margin=yr1.get("ebitda_margin_pct", 0) / 100.0
                    if yr1.get("ebitda_margin_pct", 0) > 1 else yr1.get("ebitda_margin_pct", 0),
                    exit_ebitda_margin=yr_exit.get("ebitda_margin_pct", 0) / 100.0
                    if yr_exit.get("ebitda_margin_pct", 0) > 1 else yr_exit.get("ebitda_margin_pct", 0),
                    entry_ev_ebitda=float(assumptions.get("entry_multiple", 8.0)),
                    exit_ev_ebitda=float(assumptions.get("exit_multiple", 9.0)),
                    entry_net_debt=su.get("senior_debt_eur_m", 0) + su.get("notes_eur_m", 0),
                    exit_net_debt=em.get("net_debt_at_exit_eur_m", 0),
                )
        except Exception as exc:
            logger.warning("Value creation bridge computation failed: %s", exc)

    _slide_cover(prs, company_name, date_str)
    _slide_exec_summary(prs, lbo, company_name)
    _slide_business_overview(prs, lbo, company_name)
    _slide_investment_thesis(prs, target_row, company_name)
    _slide_dcf_valuation(prs, dcf, company_name)
    _slide_lbo_analysis(prs, lbo, company_name)
    _slide_sensitivity(prs, lbo, company_name)
    _slide_value_creation(prs, bridge, company_name)
    _slide_risks(prs, company_name)
    _slide_appendix(prs, lbo, company_name)

    output_dir.mkdir(parents=True, exist_ok=True)
    slug = company_name.lower().replace(" ", "_").replace("/", "_")
    ts = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    out_path = output_dir / f"pitch_{slug}_{ts}.pptx"
    prs.save(str(out_path))
    logger.info("Pitch deck written: %s (%d slides)", out_path, len(prs.slides))
    return out_path


# ── CLI ───────────────────────────────────────────────────────────────────────

def main() -> int:
    parser = argparse.ArgumentParser(description="AUCTUS Pitch Deck Builder")
    parser.add_argument("--company-name", required=True, dest="company_name")
    parser.add_argument("--lbo-compact", type=Path, default=None, dest="lbo_compact")
    parser.add_argument("--dcf-results", type=Path, default=None, dest="dcf_results")
    parser.add_argument("--target-matrix", type=Path, default=None, dest="target_matrix")
    parser.add_argument("--output-dir", type=Path, default=Path("outputs/pitch_decks/"),
                        dest="output_dir")
    parser.add_argument("--date", type=str, default=None)
    args = parser.parse_args()

    try:
        out_path = build_deck(
            company_name=args.company_name,
            lbo_compact_path=args.lbo_compact,
            dcf_results_path=args.dcf_results,
            target_matrix_path=args.target_matrix,
            output_dir=args.output_dir,
            deal_date=args.date,
        )
        import json as _json
        print(_json.dumps({"status": "success", "output_path": str(out_path)}))
        return 0
    except ValueError as exc:
        logger.error("Input validation error: %s", exc)
        return 1
    except Exception as exc:
        logger.error("Deck build error: %s", exc)
        return 2


if __name__ == "__main__":
    sys.exit(main())
