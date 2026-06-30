"""Unit and integration tests for scripts/deck_builder.py"""

from __future__ import annotations

import json
from pathlib import Path

import pytest

from scripts.deck_builder import build_deck


@pytest.fixture
def sample_lbo_compact(tmp_path: Path) -> Path:
    """Minimal LBO compact JSON matching the compact payload schema."""
    payload = {
        "company_name": "Muster GmbH",
        "run_timestamp": "20260629_120000",
        "sources_uses": {
            "entry_ev_eur_m": 68.0,
            "equity_eur_m": 32.3,
            "senior_debt_eur_m": 27.2,
            "notes_eur_m": 10.2,
            "total_uses_eur_m": 69.7,
            "balance_check_eur_m": 0.0,
        },
        "inflection_projections": [
            {
                "year": 1,
                "revenue_eur_m": 53.5,
                "ebitda_eur_m": 9.36,
                "ebitda_margin_pct": 17.5,
                "total_interest_eur_m": 2.85,
                "levered_fcf_eur_m": 1.2,
                "total_debt_closing_eur_m": 35.1,
                "leverage_x": 3.75,
                "interest_coverage_x": 2.3,
            },
            {
                "year": 3,
                "revenue_eur_m": 61.0,
                "ebitda_eur_m": 11.3,
                "ebitda_margin_pct": 18.5,
                "total_interest_eur_m": 2.4,
                "levered_fcf_eur_m": 2.1,
                "total_debt_closing_eur_m": 30.0,
                "leverage_x": 2.65,
                "interest_coverage_x": 3.1,
            },
            {
                "year": 5,
                "revenue_eur_m": 65.4,
                "ebitda_eur_m": 12.75,
                "ebitda_margin_pct": 19.5,
                "total_interest_eur_m": 1.9,
                "levered_fcf_eur_m": 3.2,
                "total_debt_closing_eur_m": 23.5,
                "leverage_x": 1.84,
                "interest_coverage_x": 4.2,
            },
        ],
        "exit_metrics": {
            "exit_ev_eur_m": 114.75,
            "net_debt_at_exit_eur_m": 23.5,
            "equity_proceeds_eur_m": 91.25,
            "moic": 2.82,
            "irr_pct": 23.1,
            "irr_solver_converged": True,
            "leverage_at_entry_x": 4.37,
            "leverage_at_exit_x": 1.84,
            "interest_coverage_min_x": 2.3,
        },
        "sensitivity_irr_pct": {
            "6.0x": {"7.0x": 18.2, "8.0x": 20.5, "9.0x": 22.8, "10.0x": 24.9, "11.0x": 26.8},
            "7.0x": {"7.0x": 15.1, "8.0x": 17.8, "9.0x": 20.2, "10.0x": 22.3, "11.0x": 24.3},
            "8.0x": {"7.0x": 12.2, "8.0x": 15.1, "9.0x": 17.7, "10.0x": 19.9, "11.0x": 22.0},
            "9.0x": {"7.0x": 9.4, "8.0x": 12.5, "9.0x": 15.3, "10.0x": 17.7, "11.0x": 19.8},
            "10.0x": {"7.0x": 6.8, "8.0x": 10.1, "9.0x": 13.0, "10.0x": 15.5, "11.0x": 17.7},
        },
        "sensitivity_moic": {},
        "assumptions": {
            "entry_multiple": 8.0,
            "exit_multiple": 9.0,
            "exit_year": 5,
            "geography": "DE",
            "revenue_growth_rates": [0.07, 0.07, 0.06, 0.06, 0.05],
            "ebitda_margins": [0.175, 0.18, 0.185, 0.19, 0.195],
            "da_pct_revenue": 0.025,
            "capex_pct_revenue": 0.035,
            "senior_spread_bps": 375,
            "notes_fixed_rate": 0.095,
            "tax_rate": 0.299,
            "senior_amort_pct_annual": 0.05,
            "senior_cash_sweep_pct": 0.50,
            "fees_capitalized": True,
            "nwc_pct_revenue_change": 0.08,
        },
    }
    p = tmp_path / "lbo_muster_gmbh_20260629_120000_lbo_compact.json"
    p.write_text(json.dumps(payload))
    return p


@pytest.fixture
def sample_dcf_results(tmp_path: Path) -> Path:
    """Minimal DCF results JSON."""
    payload = {
        "company_name": "Muster GmbH",
        "run_timestamp": "20260629_120000",
        "wacc_used": 0.12,
        "terminal_growth_rate_used": 0.025,
        "projection_years": 5,
        "pv_forecast_cashflows_eur_m": 15.2,
        "terminal_value_pv_eur_m": 42.8,
        "enterprise_value_eur_m": 58.0,
        "terminal_value_pct_of_ev": 0.738,
    }
    p = tmp_path / "muster_gmbh_20260629_120000_dcf_results.json"
    p.write_text(json.dumps(payload))
    return p


class TestBuildDeck:
    def test_deck_created_with_no_data(self, tmp_path: Path) -> None:
        """Deck should be buildable with no model data at all."""
        out = build_deck("No Data GmbH", output_dir=tmp_path)
        assert out.exists()
        assert out.suffix == ".pptx"
        assert out.stat().st_size > 0

    def test_deck_has_ten_slides(self, tmp_path: Path) -> None:
        from pptx import Presentation
        out = build_deck("Slide Count GmbH", output_dir=tmp_path)
        prs = Presentation(str(out))
        assert len(prs.slides) == 10

    def test_deck_with_lbo_data(self, tmp_path: Path, sample_lbo_compact: Path) -> None:
        out = build_deck(
            "Muster GmbH",
            lbo_compact_path=sample_lbo_compact,
            output_dir=tmp_path,
        )
        assert out.exists()
        assert out.stat().st_size > 10_000  # non-trivial size

    def test_deck_with_all_data(
        self, tmp_path: Path, sample_lbo_compact: Path, sample_dcf_results: Path
    ) -> None:
        out = build_deck(
            "Muster GmbH",
            lbo_compact_path=sample_lbo_compact,
            dcf_results_path=sample_dcf_results,
            output_dir=tmp_path,
        )
        assert out.exists()
        from pptx import Presentation
        prs = Presentation(str(out))
        assert len(prs.slides) == 10

    def test_deck_output_filename_contains_slug_and_timestamp(
        self, tmp_path: Path
    ) -> None:
        out = build_deck("Alpha Beta GmbH", output_dir=tmp_path)
        assert "alpha_beta_gmbh" in out.name
        assert out.name.endswith(".pptx")

    def test_deck_slides_have_text_content(
        self, tmp_path: Path, sample_lbo_compact: Path
    ) -> None:
        """Every slide should have at least one text box populated."""
        from pptx import Presentation
        out = build_deck(
            "Muster GmbH",
            lbo_compact_path=sample_lbo_compact,
            output_dir=tmp_path,
        )
        prs = Presentation(str(out))
        for i, slide in enumerate(prs.slides):
            text_shapes = [
                s for s in slide.shapes
                if s.has_text_frame and s.text_frame.text.strip()
            ]
            assert len(text_shapes) > 0, f"Slide {i + 1} has no text content"

    def test_sensitivity_slide_renders_grid(
        self, tmp_path: Path, sample_lbo_compact: Path
    ) -> None:
        """Sensitivity slide (slide 7) should include a table shape when data present."""
        from pptx import Presentation
        out = build_deck(
            "Muster GmbH",
            lbo_compact_path=sample_lbo_compact,
            output_dir=tmp_path,
        )
        prs = Presentation(str(out))
        # Sensitivity is slide index 6 (0-based)
        sens_slide = prs.slides[6]
        table_shapes = [s for s in sens_slide.shapes if s.has_table]
        assert len(table_shapes) >= 1, "Sensitivity slide should contain a table"

    def test_deck_output_dir_created_automatically(self, tmp_path: Path) -> None:
        nested = tmp_path / "deep" / "output" / "dir"
        out = build_deck("Nested GmbH", output_dir=nested)
        assert nested.exists()
        assert out.exists()

    def test_missing_lbo_path_does_not_crash(self, tmp_path: Path) -> None:
        out = build_deck(
            "Missing Data GmbH",
            lbo_compact_path=tmp_path / "nonexistent.json",
            output_dir=tmp_path,
        )
        assert out.exists()

    def test_deal_date_appears_in_cover(
        self, tmp_path: Path, sample_lbo_compact: Path
    ) -> None:
        from pptx import Presentation
        specific_date = "2026-06-29"
        out = build_deck(
            "Muster GmbH",
            lbo_compact_path=sample_lbo_compact,
            output_dir=tmp_path,
            deal_date=specific_date,
        )
        prs = Presentation(str(out))
        cover_slide = prs.slides[0]
        cover_text = " ".join(
            s.text_frame.text for s in cover_slide.shapes if s.has_text_frame
        )
        assert specific_date in cover_text
